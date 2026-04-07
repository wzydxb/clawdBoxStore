#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path

PLACEHOLDER_PATTERNS = [
    '描述相关的信息以解释你的标题',
    '输入相关的描述信息以解释你的标题',
    '请在此处编辑文字',
    '公司名字',
    '主讲人：李天天',
]


def inventory_template(template_path: Path) -> dict:
    from pptx import Presentation
    prs = Presentation(str(template_path))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        shapes = []
        for shape in slide.shapes:
            text = ''
            if getattr(shape, 'has_text_frame', False) and shape.text_frame is not None:
                text = shape.text_frame.text.strip()
            shapes.append({
                'shape_id': getattr(shape, 'shape_id', None),
                'name': getattr(shape, 'name', ''),
                'left': int(getattr(shape, 'left', 0)),
                'top': int(getattr(shape, 'top', 0)),
                'width': int(getattr(shape, 'width', 0)),
                'height': int(getattr(shape, 'height', 0)),
                'text': text,
                'placeholder_hit': any(pattern in text for pattern in PLACEHOLDER_PATTERNS),
                'has_text_frame': bool(getattr(shape, 'has_text_frame', False)),
            })
        slides.append({
            'slide_number': idx,
            'shape_count': len(slide.shapes),
            'text_shapes': [s for s in shapes if s['text']],
        })
    return {
        'template': str(template_path),
        'slide_count': len(prs.slides),
        'slides': slides,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description='Inventory PPTX template structure')
    parser.add_argument('--template', required=True)
    parser.add_argument('--output', required=True)
    args = parser.parse_args()

    template_path = Path(args.template).resolve()
    output_path = Path(args.output).resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    report = inventory_template(template_path)
    output_path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + '\n', encoding='utf-8')
    print(output_path)
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
