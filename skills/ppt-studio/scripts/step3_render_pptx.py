#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import shutil
import tempfile
import zipfile
from pathlib import Path

PLACEHOLDER_PATTERNS = [
    '描述相关的信息以解释你的标题',
    '输入相关的描述信息以解释你的标题',
    '请在此处编辑文字',
    '公司名字',
    '主讲人：李天天',
]


def clear_shape_text(shape) -> None:
    if getattr(shape, 'has_text_frame', False) and shape.text_frame is not None:
        shape.text_frame.clear()
        if shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].text = ''


def set_shape_text(shape, text: str) -> None:
    if getattr(shape, 'has_text_frame', False) and shape.text_frame is not None:
        shape.text_frame.clear()
        if shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].text = text


def clean_placeholder_text(slide) -> None:
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False) and shape.text_frame is not None:
            text = shape.text_frame.text or ''
            if any(token in text for token in PLACEHOLDER_PATTERNS):
                clear_shape_text(shape)


def delete_slide(prs, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide = slides[index]
    rel_id = slide.rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide)


def market_page_payload(role: str, theme: str) -> dict:
    mapping = {
        'cover': ['市场分析报告', theme, '管理层汇报'],
        'toc': ['目录', '市场分析\n重点区域\n竞争格局', '客户画像\n核心数据\n总结建议'],
        'section-market': ['市场分析', '从市场结构和空间看机会'],
        'market-segments': ['市场机会与分层', '低端市场\n中端市场\n高端市场', '聚焦最具复制性的中端场景'],
        'region-data': ['重点城市与区域', '核心城市 / 区域优先级', '用真实区域数据替换示意项'],
        'competition': ['竞争格局', '主要竞争者 / 对标对象', '用差异化切入替换模板说明'],
        'section-customer': ['客户分析', '从客户结构与画像理解需求'],
        'persona': ['用户画像', '高价值客户特征\n消费动机\n渠道偏好', '画像要服务后续动作'],
        'core-data': ['核心数据判断', '市场规模 / 趋势 / 核心指标', '示意数据必须明确标注'],
        'recommendation': ['总结与建议', '聚焦高确定性场景\n优先验证可复制路径\n控制试点成本', '建议直接可执行'],
        'closing': ['感谢观看'],
    }
    return {'texts': mapping.get(role, ['感谢观看'])}


def training_page_payload(role: str, theme: str) -> dict:
    mapping = {
        'cover': [theme, '员工培训指南', '内部培训 / 制度宣导'],
        'toc': ['目录', '培训目的\n关键知识点\n流程规范', '案例示例\n风险提醒\n培训总结'],
        'training-goal': ['培训目的', '统一认知 / 建立规则 / 降低风险', '为什么要做这次培训'],
        'key-point-1': ['关键知识点一', '核心概念\n行为要求\n常见误区', '先讲规则，再讲案例'],
        'key-point-2': ['关键知识点二', '重点规则\n执行边界\n检查方式', '明确什么能做、什么不能做'],
        'process': ['流程与规范', '步骤一 → 步骤二 → 步骤三 → 步骤四', '按流程执行，减少偏差'],
        'example': ['案例示例', '正确示例 / 错误示例 / 改进方式', '用真实案例帮助理解'],
        'risk': ['风险提醒', '高风险行为\n必须避免的错误', '出现问题时的处理方式'],
        'summary': ['培训总结', '回顾关键点 / 明确执行要求', '培训后立即可执行'],
        'closing': ['感谢聆听'],
    }
    return {'texts': mapping.get(role, ['感谢聆听'])}


def fill_text_shapes(slide, texts: list[str]) -> None:
    text_shapes = [shape for shape in slide.shapes if getattr(shape, 'has_text_frame', False) and shape.text_frame is not None]
    for idx, text in enumerate(texts):
        if idx < len(text_shapes):
            set_shape_text(text_shapes[idx], text)


def scrub_placeholders_in_pptx(path: Path) -> None:
    replacements = [
        '描述相关的信息以解释你的标题。现在就开始打字吧。写任何你想表达的内容',
        '输入相关的描述信息以解释你的标题。',
        '请在此处编辑文字',
        '公司名字',
        '主讲人：李天天',
        '资深研究员，创业公司CEO',
    ]
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        with zipfile.ZipFile(path, 'r') as zin:
            zin.extractall(tmp_path)
        for xml_path in (tmp_path / 'ppt' / 'slides').glob('slide*.xml'):
            text = xml_path.read_text(encoding='utf-8', errors='ignore')
            for old in replacements:
                text = text.replace(old, '')
            xml_path.write_text(text, encoding='utf-8')
        backup = path.with_suffix('.bak')
        shutil.copy2(path, backup)
        with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for file in tmp_path.rglob('*'):
                if file.is_file():
                    zout.write(file, file.relative_to(tmp_path))
        backup.unlink(missing_ok=True)


def build_pptx(plan: dict, output: Path, selected_template: dict | None, template_dir: Path, mapping_dir: Path) -> str:
    from pptx import Presentation
    template_path = template_dir / selected_template['file'] if selected_template else None
    mode = 'scratch-mvp-python-pptx'
    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
        mode = 'template-fill-python-pptx'
    else:
        prs = Presentation()

    template_id = selected_template.get('id') if selected_template else None
    mapping_path = mapping_dir / f'{template_id}.json' if template_id else None
    mapping = json.loads(mapping_path.read_text(encoding='utf-8')) if mapping_path and mapping_path.exists() else None
    theme = plan.get('theme', '未命名主题')

    if mapping and prs.slides:
        for page in mapping['pages']:
            idx = page['slide_number'] - 1
            if idx >= len(prs.slides):
                continue
            slide = prs.slides[idx]
            clean_placeholder_text(slide)
            role = page['role']
            payload = market_page_payload(role, theme) if template_id == 'market-analysis-report' else training_page_payload(role, theme)
            fill_text_shapes(slide, payload['texts'])
        while len(prs.slides) > len(mapping['pages']):
            delete_slide(prs, len(prs.slides) - 1)
    else:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        text_box = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        text_box.text_frame.text = theme

    prs.save(output)
    scrub_placeholders_in_pptx(output)
    return mode


def main() -> int:
    parser = argparse.ArgumentParser(description='Render final PPTX for ppt-studio')
    parser.add_argument('--plan', required=True)
    parser.add_argument('--specs', required=True)
    parser.add_argument('--assets', required=True)
    parser.add_argument('--style-resolution')
    parser.add_argument('--output', required=True)
    args = parser.parse_args()

    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    plan = json.loads(Path(args.plan).read_text(encoding='utf-8'))
    style_resolution = json.loads(Path(args.style_resolution).read_text(encoding='utf-8')) if args.style_resolution else {}
    selected_template = style_resolution.get('selected_template')
    template_dir = Path(__file__).resolve().parents[1] / 'assets' / 'templates'
    mapping_dir = template_dir / 'mappings'
    mode = build_pptx(plan, out_path, selected_template, template_dir, mapping_dir)
    report = {
        'plan': args.plan,
        'specs': args.specs,
        'assets': args.assets,
        'style_resolution': args.style_resolution,
        'selected_template': selected_template,
        'output': str(out_path),
        'mode': mode,
        'slide_count': len(plan.get('pages', [])),
    }
    (out_path.parent / 'render_report.json').write_text(json.dumps(report, ensure_ascii=False, indent=2) + '\n', encoding='utf-8')
    print(out_path)
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
