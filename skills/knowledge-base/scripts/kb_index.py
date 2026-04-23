#!/usr/bin/env python3
"""
kb_index.py - 知识库索引管理

用法：
  python3 kb_index.py <mount_path>              # 增量扫描建索引
  python3 kb_index.py <mount_path> --update <file_path>  # 单文件更新并 diff
  python3 kb_index.py <mount_path> --init       # 初始化（首次建库）
"""

import sys
import os
import sqlite3
import difflib
import argparse
from pathlib import Path
from datetime import datetime

SKIP_DIRS = {".hermes-index", ".DS_Store", ".Spotlight-V100", ".Trashes"}
SENSITIVE_KEYWORDS = {"密码", "私钥", "身份证", "password", "secret", "credential"}
MAX_CONTENT = 50000  # 每个文件最多提取字符数


# ── 内容提取 ──────────────────────────────────────────────

def extract_content(path: str) -> str:
    ext = Path(path).suffix.lower()
    try:
        if ext == ".pdf":
            import pymupdf
            doc = pymupdf.open(path)
            return "\n".join(p.get_text() for p in doc)[:MAX_CONTENT]

        elif ext == ".docx":
            from docx import Document
            return "\n".join(p.text for p in Document(path).paragraphs)[:MAX_CONTENT]

        elif ext == ".xlsx":
            import pandas as pd
            return pd.read_excel(path).to_string()[:MAX_CONTENT]

        elif ext == ".csv":
            import pandas as pd
            return pd.read_csv(path, encoding_errors="ignore").to_string()[:MAX_CONTENT]

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            texts = [s.text for sl in prs.slides
                     for s in sl.shapes if s.has_text_frame]
            return "\n".join(texts)[:MAX_CONTENT]

        elif ext in (".txt", ".md", ".json", ".yaml", ".yml",
                     ".py", ".js", ".ts", ".sh", ".html", ".xml"):
            return open(path, errors="ignore").read()[:MAX_CONTENT]

        else:
            return ""  # 图片/视频只存元数据

    except Exception as e:
        return f"[提取失败: {e}]"


def is_sensitive(fname: str) -> bool:
    return any(k in fname.lower() for k in SENSITIVE_KEYWORDS)


# ── 数据库初始化 ──────────────────────────────────────────

def init_db(mount: str) -> sqlite3.Connection:
    index_dir = os.path.join(mount, ".hermes-index")
    os.makedirs(os.path.join(index_dir, "wiki", "concepts"), exist_ok=True)
    os.makedirs(os.path.join(index_dir, "wiki", "insights"), exist_ok=True)

    db = sqlite3.connect(os.path.join(index_dir, "index.db"))
    db.execute("""
        CREATE VIRTUAL TABLE IF NOT EXISTS files USING fts5(
            path, name, content, tags, summary,
            tokenize='trigram'
        )
    """)
    db.execute("""
        CREATE TABLE IF NOT EXISTS meta (
            path TEXT PRIMARY KEY,
            mtime REAL,
            size  INTEGER,
            filetype TEXT,
            indexed_at TEXT
        )
    """)
    db.commit()
    return db


# ── 增量扫描 ──────────────────────────────────────────────

def scan(mount: str):
    db = init_db(mount)
    added = updated = skipped = 0

    for root, dirs, files in os.walk(mount):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
        for fname in files:
            if fname.startswith("."):
                continue
            fpath = os.path.join(root, fname)
            try:
                stat = os.stat(fpath)
            except OSError:
                continue

            row = db.execute("SELECT mtime FROM meta WHERE path=?",
                             (fpath,)).fetchone()
            if row and row[0] == stat.st_mtime:
                skipped += 1
                continue

            content = "[敏感文件，内容跳过]" if is_sensitive(fname) \
                      else extract_content(fpath)

            db.execute("DELETE FROM files WHERE path=?", (fpath,))
            db.execute("INSERT INTO files VALUES (?,?,?,?,?)",
                       (fpath, fname, content, "", ""))
            db.execute("""
                INSERT OR REPLACE INTO meta VALUES (?,?,?,?,?)
            """, (fpath, stat.st_mtime, stat.st_size,
                  Path(fpath).suffix, datetime.now().isoformat()))

            if row:
                updated += 1
            else:
                added += 1

    db.commit()
    print(f"扫描完成：新增 {added} | 更新 {updated} | 跳过 {skipped}")

    # 更新 wiki log
    _append_log(mount, f"增量扫描：新增 {added}，更新 {updated}")


# ── 单文件更新 + diff ─────────────────────────────────────

def update_file(mount: str, file_path: str):
    db = init_db(mount)

    # 1. 取旧内容
    row = db.execute("SELECT content FROM files WHERE path=?",
                     (file_path,)).fetchone()
    old_content = row[0] if row else ""

    # 2. 提取新内容
    new_content = "[敏感文件，内容跳过]" if is_sensitive(Path(file_path).name) \
                  else extract_content(file_path)

    # 3. diff（先 diff 再覆盖）
    diff_lines = list(difflib.unified_diff(
        old_content.splitlines(),
        new_content.splitlines(),
        fromfile="旧版本",
        tofile="新版本",
        lineterm="",
        n=3
    ))

    # 4. 更新索引
    fname = Path(file_path).name
    stat = os.stat(file_path)
    db.execute("DELETE FROM files WHERE path=?", (file_path,))
    db.execute("INSERT INTO files VALUES (?,?,?,?,?)",
               (file_path, fname, new_content, "", ""))
    db.execute("INSERT OR REPLACE INTO meta VALUES (?,?,?,?,?)",
               (file_path, stat.st_mtime, stat.st_size,
                Path(file_path).suffix, datetime.now().isoformat()))
    db.commit()

    # 5. 输出 diff 供 Agent 解读
    print(f"\n📄 {fname} 索引已更新\n")
    if diff_lines:
        print("--- 变更明细 ---")
        for line in diff_lines[:100]:  # 最多展示100行
            print(line)
        if len(diff_lines) > 100:
            print(f"... 共 {len(diff_lines)} 行变更，已截断")
    else:
        print("（内容无变化）")

    # 6. 更新 wiki log
    _append_log(mount, f"文件更新：{fname}")


# ── Wiki log ─────────────────────────────────────────────

def _append_log(mount: str, message: str):
    log_path = os.path.join(mount, ".hermes-index", "wiki", "log.md")
    ts = datetime.now().strftime("%Y-%m-%d %H:%M")
    with open(log_path, "a") as f:
        f.write(f"- {ts}：{message}\n")


# ── 入口 ─────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="知识库索引管理")
    parser.add_argument("mount", help="网盘挂载路径，如 /Volumes/uploads")
    parser.add_argument("--update", metavar="FILE", help="更新单个文件并输出 diff")
    parser.add_argument("--init", action="store_true", help="仅初始化数据库结构")
    args = parser.parse_args()

    if not os.path.exists(args.mount):
        print(f"错误：挂载路径不存在 {args.mount}")
        sys.exit(1)

    if args.init:
        init_db(args.mount)
        print(f"初始化完成：{args.mount}/.hermes-index/")
    elif args.update:
        update_file(args.mount, args.update)
    else:
        scan(args.mount)
